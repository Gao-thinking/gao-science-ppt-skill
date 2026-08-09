#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gao-science-ppt-skill: PPTX 内容诊断脚本
提取每页标题(大字文本启发式)/全部文本(字号/字体/粗体)/图片(宽高/比例/压缩状态)/文字量
输出 JSON 诊断报告 + 可读摘要, 供 SKILL.md §3 Step2 贝叶斯分级使用。

用法:
  python3 pptx_analyze.py -i deck.pptx -o /tmp/ppt-diag/
"""
import argparse
import json
import os
import sys
from collections import Counter

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("缺少 python-pptx，请先: python3 -m pip install --user python-pptx")


# 字号分档: 用于"标题 vs 正文"启发式
SIZE_BANDS = [
    (40, "cover"),      # 封面/大标题
    (32, "section"),    # 节标题
    (24, "title"),      # 页标题
    (18, "body"),       # 正文
    (0, "note"),        # 注释
]


def band_of(size_pt):
    for th, name in SIZE_BANDS:
        if size_pt >= th:
            return name
    return "note"


def iter_texts(shape):
    """遍历 shape 内所有段落文字, 产出 (text, size_pt, font_name, bold)。"""
    out = []
    if not getattr(shape, "has_text_frame", False):
        return out
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if not t:
                continue
            sz = run.font.size.pt if run.font.size else None
            fname = run.font.name
            bold = bool(run.font.bold)
            out.append({"text": t, "size": sz, "font": fname, "bold": bold})
    return out


def shape_images(shape):
    """收集 shape 图片: 宽高(英寸)/比例/压缩状态。"""
    out = []
    if shape.shape_type == 13:  # PICTURE
        w_in = shape.width / 914400.0
        h_in = shape.height / 914400.0
        ratio = round(w_in / h_in, 2) if h_in else 0
        # 压缩状态: 显示尺寸 vs 原图尺寸
        compressed = None
        try:
            img = shape.image
            ow_in = img.size[0] / 914400.0 if False else None  # placeholder
            ow = img.size[0]  # px
            oh = img.size[1]
            if ow and oh:
                disp_px_w = shape.width / 914400.0 * 96
                disp_px_h = shape.height / 914400.0 * 96
                compressed = (disp_px_w > ow * 1.5) or (disp_px_h > oh * 1.5)
        except Exception:
            pass
        out.append({"w_in": round(w_in, 2), "h_in": round(h_in, 2),
                    "ratio": ratio, "compressed": compressed})
    return out


def pick_title(texts):
    """大字文本启发式: 取字号最大且字数 <= 30 的文本为标题。"""
    if not texts:
        return None
    t = max(texts, key=lambda x: x["size"] or 0)
    if t["size"] and t["size"] >= 24 and len(t["text"]) <= 30:
        return t
    return None


def analyze(path):
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        images = []
        for shape in slide.shapes:
            texts.extend(iter_texts(shape))
            images.extend(shape_images(shape))
        # 文本文字量(去空格)
        word_count = sum(len(t["text"].replace(" ", "")) for t in texts)
        title = pick_title(texts)
        # 页类型启发式
        page_type = "content"
        if idx == 1:
            page_type = "cover"
        elif word_count <= 10:
            page_type = "section"
        elif any("练习" in t["text"] or "习题" in t["text"] for t in texts):
            page_type = "exercise"
        elif any("小结" in t["text"] or "总结" in t["text"] for t in texts):
            page_type = "summary"
        slides.append({
            "index": idx,
            "title": title["text"] if title else None,
            "texts": texts,
            "images": images,
            "word_count": word_count,
            "page_type": page_type,
        })
    return slides


def summarize(slides):
    """可读摘要: 字体种类/字号档位/图片比例分布/超载页。"""
    lines = [f"共 {len(slides)} 页"]
    all_fonts = Counter()
    all_sizes = Counter()
    ratios = Counter()
    over = []
    title_missing = []
    for s in slides:
        for t in s["texts"]:
            if t["font"]:
                all_fonts[t["font"]] += 1
            if t["size"]:
                all_sizes[round(t["size"])] += 1
        for im in s["images"]:
            ratios[im["ratio"]] += 1
        if s["word_count"] > 150:
            over.append((s["index"], s["word_count"]))
        if not s["title"]:
            title_missing.append(s["index"])
    lines.append(f"字体种类 {len(all_fonts)}: {dict(all_fonts.most_common(5))}")
    lines.append(f"字号档位 {len(all_sizes)}: {dict(all_sizes.most_common(6))}")
    lines.append(f"图片比例 {len(ratios)}: {dict(ratios)}")
    if over:
        lines.append(f"超载页(>150字): {over}")
    if title_missing:
        lines.append(f"无标题页: {title_missing}")
    return "\n".join(lines)


def main():
    ap = argparse.ArgumentParser(description="PPTX 课件内容诊断")
    ap.add_argument("-i", "--input", required=True, help="输入 .pptx")
    ap.add_argument("-o", "--outdir", default="/tmp/ppt-diag/", help="输出目录")
    args = ap.parse_args()

    if not args.input.endswith(".pptx"):
        sys.exit("仅支持 .pptx（.ppt 请先用 Office/LibreOffice 另存为 .pptx）")

    os.makedirs(args.outdir, exist_ok=True)
    base = os.path.splitext(os.path.basename(args.input))[0]
    slides = analyze(args.input)
    report = {"source": args.input, "pages": len(slides), "slides": slides}
    json_path = os.path.join(args.outdir, f"{base}-诊断.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=1)
    print(summarize(slides))
    print(f"\n→ {json_path}")


if __name__ == "__main__":
    main()
