#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gao-science-ppt-skill: 兜底工具——逐页打印文本框/图片的 名称、类型、位置、文本，
用于 §3A 里 toc/box/para_replace 需要按形状名定位时的取形帮助（缺省启发式选错时）。

用法:
  python3 pptx_dump_shapes.py -i deck.pptx [--pages 1,2,5]
"""
import argparse
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("缺少 python-pptx，请先: python3 -m pip install --user python-pptx")


def main():
    ap = argparse.ArgumentParser(description="逐页 dump 形状名/类型/位置/文本")
    ap.add_argument("-i", "--input", required=True, help="输入 .pptx")
    ap.add_argument("--pages", default="", help="逗号分隔页号，缺省=全部")
    args = ap.parse_args()
    prs = Presentation(args.input)
    pages = {int(p) for p in args.pages.split(",")} if args.pages else None
    for idx, slide in enumerate(prs.slides, 1):
        if pages and idx not in pages:
            continue
        print(f"\n===== slide {idx} =====")
        for sh in slide.shapes:
            pos = ""
            if sh.left is not None:
                pos = (f"L{sh.left/914400:.2f} T{sh.top/914400:.2f} "
                       f"W{sh.width/914400:.2f} H{sh.height/914400:.2f}")
            txt = ""
            if getattr(sh, "has_text_frame", False):
                txt = " | ".join(
                    "".join(r.text for r in p.runs)
                    for p in sh.text_frame.paragraphs)
            print(f"  [{sh.shape_type}] {pos} name={sh.name!r} :: {txt[:60]}")


if __name__ == "__main__":
    main()