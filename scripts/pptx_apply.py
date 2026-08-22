#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gao-science-ppt-skill: PPTX 排版统一应用脚本
按方案 JSON 执行: 标题重写 / 字体统一(中英分离) / 字号统一 / 图片统一(等宽或裁剪) / 教学设计页填充。
只改排版不改文字事实。默认输出新文件, 不覆盖原件。

用法:
  python3 pptx_apply.py -i deck.pptx -o deck-优化.pptx --plan plan.json [--append-teaching teaching.json]

方案 JSON 结构:
{
  "titles": {"<页号>": "<新标题>"},                 # 页号 = 诊断 JSON 的 slide.index
  "fonts": {"latin": "Calibri", "ea": "微软雅黑"},  # latin=西文/数字, ea=中文
  "size_map": {"18": 20, "16": 20},                 # 旧字号pt → 新字号pt
  "images": {"mode": "equal_width",                 # equal_width | crop_to_ratio | none
             "width_in": 5.5, "ratio": "4:3",
             "align": "center", "border": true},
  "teaching": {...}                                  # 或单独 --append-teaching
}
"""
import argparse
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
except ImportError:
    sys.exit("缺少 python-pptx，请先: python3 -m pip install --user python-pptx")

GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DARK = RGBColor(0x33, 0x33, 0x33)


def set_run_fonts(run, latin=None, ea=None):
    """同时设置 latin(西文) 与 ea(中文) 字体，按 OOXML schema 顺序插入
    (latin → ea → cs → sym)，避免 rPr 子元素乱序导致 PowerPoint 报 repair。"""
    if latin:
        run.font.name = latin
    if ea:
        rPr = run._r.get_or_add_rPr()
        # 先移除旧 ea/cs，避免重复
        for tag in (qn("a:ea"), qn("a:cs")):
            el = rPr.find(tag)
            if el is not None:
                rPr.remove(el)
        ea_el = rPr.makeelement(qn("a:ea"), {"typeface": ea})
        cs_el = rPr.makeelement(qn("a:cs"), {"typeface": ea})
        latin_el = rPr.find(qn("a:latin"))
        if latin_el is not None:
            # latin → ea → cs
            latin_el.addnext(cs_el)
            latin_el.addnext(ea_el)
        else:
            # 无 latin：插到 sym/hlinkClick/hlinkMouseOver/rtl/extLst 之前
            anchor = None
            for tag in (qn("a:sym"), qn("a:hlinkClick"), qn("a:hlinkMouseOver"),
                        qn("a:rtl"), qn("a:extLst")):
                el = rPr.find(tag)
                if el is not None:
                    anchor = el
                    break
            if anchor is not None:
                anchor.addprevious(cs_el)
                anchor.addprevious(ea_el)
            else:
                rPr.append(ea_el)
                rPr.append(cs_el)


def iter_runs(shape):
    """递归收集 run：支持 GROUP 组合内文字与表格单元格文字。
    用于字体/字号统一，保证不遗漏组合与表格内的 run。"""
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            yield from iter_runs(child)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        yield run
        return
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run


def replace_text(shape, new_text):
    """把 shape 文本替换为新文本(保留首 run 格式)。"""
    tf = shape.text_frame
    runs = tf.paragraphs[0].runs
    if runs:
        runs[0].text = new_text
        for extra in runs[1:]:
            extra.text = ""
    else:
        tf.text = new_text
    # 清掉其余段落的文字
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def apply_titles(prs, title_map):
    """按页号替换标题: 找该页字号最大的文本 shape 替换。"""
    for idx, slide in enumerate(prs.slides, 1):
        new_title = title_map.get(str(idx))
        if not new_title:
            continue
        best = None
        best_sz = -1
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            sz = 0
            for run in iter_runs(shape):
                if run.font.size and run.font.size.pt > sz:
                    sz = run.font.size.pt
            if sz > best_sz:
                best, best_sz = shape, sz
        if best is not None:
            replace_text(best, new_title)


def apply_fonts(prs, fonts):
    latin = fonts.get("latin")
    ea = fonts.get("ea")
    for slide in prs.slides:
        for shape in slide.shapes:
            for run in iter_runs(shape):
                set_run_fonts(run, latin=latin, ea=ea)


def apply_sizes(prs, size_map):
    smap = {float(k): float(v) for k, v in size_map.items()}
    for slide in prs.slides:
        for shape in slide.shapes:
            for run in iter_runs(shape):
                if run.font.size and run.font.size.pt in smap:
                    run.font.size = Pt(smap[run.font.size.pt])


def apply_table_replace(prs, spec):
    """按页号+行列替换表格单元格文本（保留单元格格式）。
    spec: {"<页号>": [{"row": r, "col": c, "text": "..."}]}
    或 items 形式保持「序号金色 + 内容白色」结构：
    {"<页号>": [{"row": r, "col": c, "items": [{"num": "1.", "text": "..."}, {"num": "2.", "text": "..."}]}]}
    items 时：每段首 run 保留序号（金色），第二 run 写内容（保留原白色），多余 run/段落清空。"""
    for slide_idx, rules in spec.items():
        slide = prs.slides[int(slide_idx) - 1]
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            tbl = shape.table
            for rule in rules:
                r, c = rule["row"], rule["col"]
                if r >= len(tbl.rows) or c >= len(tbl.columns):
                    continue
                cell = tbl.cell(r, c)
                tf = cell.text_frame
                items = rule.get("items")
                if items:
                    for pi, item in enumerate(items):
                        para = tf.paragraphs[pi] if pi < len(tf.paragraphs) else tf.add_paragraph()
                        runs = para.runs
                        if not runs:
                            runs = [para.add_run()]
                        if isinstance(item, dict) and "num" in item:
                            # 序号 run：保留原格式（金色）
                            runs[0].text = item["num"]
                            content = item["text"]
                            if len(runs) > 1:
                                runs[1].text = content
                                for extra in runs[2:]:
                                    extra.text = ""
                            else:
                                r2 = para.add_run()
                                r2.text = content
                        else:
                            txt = item["text"] if isinstance(item, dict) else item
                            runs[0].text = txt
                            for extra in runs[1:]:
                                extra.text = ""
                    # 清掉多余段落
                    for extra_p in tf.paragraphs[len(items):]:
                        for run in extra_p.runs:
                            run.text = ""
                else:
                    text = rule["text"]
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = text
                        for para in tf.paragraphs:
                            for run in para.runs[1:]:
                                run.text = ""
                    else:
                        tf.text = text


def apply_text_replace(prs, spec):
    """按页号替换段落文本（段落级拼接匹配，兼容一段拆多 run）。
    spec: {"<页号>": [{"old": ..., "new": ..., "mode": "exact|contains|startswith"}]}
    保留首 run 格式，其余 run 清空。"""
    for slide_idx, rules in spec.items():
        slide = prs.slides[int(slide_idx) - 1]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                for rule in rules:
                    old, new = rule["old"], rule["new"]
                    mode = rule.get("mode", "contains")
                    if mode == "exact":
                        hit = full == old
                    elif mode == "startswith":
                        hit = full.startswith(old)
                    else:
                        hit = old in full
                    if hit:
                        if mode == "exact":
                            full = new
                        elif mode == "startswith":
                            full = new + full[len(old):]
                        else:
                            full = full.replace(old, new, 1)
                        if para.runs:
                            para.runs[0].text = full
                            for r in para.runs[1:]:
                                r.text = ""
                        break


def apply_replace_images(prs, spec):
    """按页号替换图片内容（保持原位置/尺寸，新图按原图框比例中心裁剪防拉伸）。
    spec: {"<页号>": [{"pic": <该页第几张图,0-based>, "path": "新图路径"}]}"""
    try:
        from PIL import Image, ImageOps
    except ImportError:
        print("[warn] 缺 PIL，图片替换跳过（python3 -m pip install --user Pillow）")
        return
    import io
    for slide_idx, rules in spec.items():
        slide = prs.slides[int(slide_idx) - 1]
        pics = _collect_pics(list(slide.shapes), [])
        for rule in rules:
            n, path = rule["pic"], rule["path"]
            if n >= len(pics) or not os.path.exists(path):
                continue
            pic = pics[n]
            try:
                # 目标比例 = 原图框比例
                target_ratio = (pic.width / pic.height) if pic.height else 1.0
                im = Image.open(path)
                im = ImageOps.exif_transpose(im).convert("RGB")
                w, h = im.size
                cur_ratio = w / h
                if cur_ratio > target_ratio + 0.01:   # 太宽 → 裁左右
                    nw = int(h * target_ratio)
                    x0 = (w - nw) // 2
                    im = im.crop((x0, 0, x0 + nw, h))
                elif cur_ratio < target_ratio - 0.01:  # 太高 → 裁上下
                    nh = int(w / target_ratio)
                    y0 = (h - nh) // 2
                    im = im.crop((0, y0, w, y0 + nh))
                buf = io.BytesIO()
                im.save(buf, format="JPEG", quality=88)
                tmp_dir = "/tmp/ppt-assets"
                os.makedirs(tmp_dir, exist_ok=True)
                tmp = os.path.join(tmp_dir, "_tmp_repl.jpg")
                with open(tmp, "wb") as f:
                    f.write(buf.getvalue())
                image_part, rId = slide.part.get_or_add_image_part(tmp)
                blip = pic._element.blipFill.blip
                blip.set(qn("r:embed"), rId)
            except Exception as e:
                print(f"[warn] 页{slide_idx} 图{n} 替换失败: {e}")


def _collect_pics(shapes, out):
    """递归收集所有 PICTURE（含 GROUP 内嵌），保持文档顺序。"""
    for s in shapes:
        if s.shape_type == 13:
            out.append(s)
        elif s.shape_type == 6:  # GROUP
            _collect_pics(s.shapes, out)
    return out


def set_round_rect(pic, adj="val 6060", shadow=True):
    """把图片几何改为圆角矩形（参考页2 左图样式: roundRect + adj + 阴影）。"""
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:ln", "a:effectLst"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
    avLst = geom.makeelement(qn("a:avLst"), {})
    gd = avLst.makeelement(qn("a:gd"), {"name": "adj", "fmla": adj})
    avLst.append(gd)
    geom.append(avLst)
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.addnext(geom)
    else:
        spPr.insert(0, geom)
    if shadow:
        eff = spPr.makeelement(qn("a:effectLst"), {})
        shdw = eff.makeelement(qn("a:outerShdw"), {
            "blurRad": "292100", "dist": "139700",
            "dir": "2700000", "algn": "tl", "rotWithShape": "0"})
        clr = shdw.makeelement(qn("a:srgbClr"), {"val": "333333"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": "65000"})
        clr.append(alpha)
        shdw.append(clr)
        eff.append(shdw)
        spPr.append(eff)


def apply_round_rect(prs, cfg):
    """正文图片统一圆角矩形样式——**仅针对原本就是方形(rect)的图片**，
    原本是圆形(ellipse)/已圆角(roundRect)的保持原样，避免把圆图改方。
    cfg: {"pages": [1-based] 缺省=全部, "skip": [页号], "adj": "val 6060", "shadow": true}"""
    pages = cfg.get("pages")
    skip = set(cfg.get("skip", []))
    adj = cfg.get("adj", "val 6060")
    shadow = cfg.get("shadow", True)
    for idx, slide in enumerate(prs.slides, 1):
        if pages and idx not in pages:
            continue
        if idx in skip:
            continue
        for pic in _collect_pics(list(slide.shapes), []):
            spPr = pic._element.spPr
            prst = spPr.find(qn("a:prstGeom"))
            geom = prst.get("prst") if prst is not None else None
            if geom == "rect":  # 只有方形图片才改圆角
                set_round_rect(pic, adj, shadow)


def apply_layout_images(prs, spec):
    """按页设置图片位置/尺寸（用于封面四图适配、目录左图调整）。
    spec: {"<页号>": [{"pic": n, "left_in": x, "top_in": y, "width_in": w, "height_in": h}]}
    只改出现的字段，单位英寸。"""
    for slide_idx, rules in spec.items():
        slide = prs.slides[int(slide_idx) - 1]
        pics = _collect_pics(list(slide.shapes), [])
        for rule in rules:
            n = rule["pic"]
            if n >= len(pics):
                continue
            p = pics[n]
            if "left_in" in rule:
                p.left = Inches(rule["left_in"])
            if "top_in" in rule:
                p.top = Inches(rule["top_in"])
            if "width_in" in rule:
                p.width = Inches(rule["width_in"])
            if "height_in" in rule:
                p.height = Inches(rule["height_in"])


def apply_images(prs, cfg):
    mode = cfg.get("mode", "equal_width")
    if mode == "none":
        return
    width_in = cfg.get("width_in", 5.5)
    pages = cfg.get("pages")  # 限定页号列表(1-based)，缺省=全部
    ratio_spec = cfg.get("ratio", "4:3")
    try:
        rw, rh = (float(x) for x in ratio_spec.split(":"))
        target_ratio = rw / rh
    except Exception:
        target_ratio = 4 / 3
    slide_w = prs.slide_width  # EMU
    for idx, slide in enumerate(prs.slides, 1):
        if pages and idx not in pages:
            continue
        pics = [s for s in slide.shapes if s.shape_type == 13]
        if not pics:
            continue
        max_w = max(p.width for p in pics)
        width = min(Inches(width_in), max_w)
        for p in pics:
            orig_ratio = (p.width / p.height) if p.height else target_ratio
            p.width = width
            if mode == "equal_width":
                p.height = int(width / orig_ratio)
            else:  # crop_to_ratio: 等宽 + 中心裁剪到目标比例
                p.height = int(width / target_ratio)
                _set_crop(p, orig_ratio, target_ratio)
            if cfg.get("align") == "center":
                p.left = int((slide_w - width) / 2)
            if cfg.get("border"):
                p.line.color.rgb = GRAY
                p.line.width = Pt(1)


def _set_crop(pic, orig_ratio, target_ratio):
    """中心裁剪: 用 crop 属性切到目标比例(不拉伸)。"""
    if orig_ratio > target_ratio:  # 太宽 → 裁左右
        total = 1 - target_ratio / orig_ratio
        pic.crop_left = total / 2
        pic.crop_right = total / 2
        pic.crop_top = 0
        pic.crop_bottom = 0
    else:  # 太高 → 裁上下
        total = 1 - orig_ratio / target_ratio
        pic.crop_top = total / 2
        pic.crop_bottom = total / 2
        pic.crop_left = 0
        pic.crop_right = 0


def add_teaching_slide(prs, teaching):
    """在末尾插入「教学设计」页。"""
    layout = None
    for l in prs.slide_layouts:
        if l.name and ("blank" in l.name.lower() or "空白" in l.name):
            layout = l
            break
    if layout is None:
        layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    sw = prs.slide_width
    # 标题
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                  sw - Inches(1.0), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = teaching.get("title", "教学设计")
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK
    # 内容
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.3),
                                    sw - Inches(1.6), Inches(5.0))
    bf = body.text_frame
    bf.word_wrap = True
    sections = [
        ("教学目标", teaching.get("goals", [])),
        ("教学重难点", teaching.get("key_points", [])
         + teaching.get("difficult_points", [])),
        ("教学环节", teaching.get("stages", [])),
    ]
    first = True
    for label, items in sections:
        if not items:
            continue
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"【{label}】"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = DARK
        for item in items:
            p = bf.add_paragraph()
            run = p.add_run()
            run.text = f"· {item}"
            run.font.size = Pt(18)


# ---------------------------------------------------------------------------
# §3.5 四类可复用改造的新 plan 键：toc / para_replace / box
# ---------------------------------------------------------------------------

def _pick_text_shape(slide):
    """启发式找一个文本框：优先取 run 数最多（最可能是有编号列表/正文框）。"""
    best, bestn = None, -1
    for s in slide.shapes:
        if not getattr(s, "has_text_frame", False):
            continue
        n = sum(1 for p in s.text_frame.paragraphs for _ in p.runs)
        if n > bestn:
            best, bestn = s, n
    return best


def apply_toc(prs, spec):
    """重建/重排目录条目。保留每条目首 run 格式；条目结构：
    spec: {"<页号>": {"shape": "形状名(可选)", "items": [{"num": "1.", "text": "放射性"}, ...]}}
    末尾想保留「课堂总结/练习与应用/提升训练」时，把这三条也排进 items 且序号顺延即可。
    - 若原段落含≥2个 run（常为「序号金色 + 内容白色」），num 写入 run0、text 写入 run1，保持颜色；
    - 其余情况则整段写入 run0（去掉多余 run）。
    """
    for sidx, cfg in spec.items():
        slide = prs.slides[int(sidx) - 1]
        shape = None
        if cfg.get("shape"):
            shape = next((s for s in slide.shapes if s.name == cfg["shape"]), None)
        if shape is None:
            shape = _pick_text_shape(slide)
        if shape is None or not getattr(shape, "has_text_frame", False):
            print(f"[warn] toc 页{sidx} 无目标文本框，跳过")
            continue
        tf = shape.text_frame
        paras = tf.paragraphs
        items = cfg["items"]
        for p in paras:
            for r in p.runs:
                r.text = ""
        for i, it in enumerate(items):
            num, text = "", ""
            if isinstance(it, dict):
                num, text = it.get("num", ""), it.get("text", "")
            else:
                text = str(it)
            if i < len(paras):
                p = paras[i]
                runs = p.runs
                # 定位最后一个非空 run（真正承载内容的 run），
                # 避免把内容写进中间的空格/装饰 run 而继承序号颜色（如金色）
                last_ne = -1
                for j, r in enumerate(runs):
                    if r.text.strip():
                        last_ne = j
                if runs and last_ne >= 1:
                    runs[0].text = num
                    for j, r in enumerate(runs):
                        if j == 0:
                            continue
                        r.text = (" " + text) if j == last_ne else ""
                elif runs:
                    runs[0].text = (num + " " + text).strip()
                    for extra in runs[1:]:
                        extra.text = ""
                else:
                    para = p
                    if not para.runs:
                        para.add_run()
                    para.runs[0].text = (num + " " + text).strip()
            else:
                p = tf.add_paragraph()
                p.add_run().text = (num + " " + text).strip()
        for p in paras[len(items):]:
            for r in p.runs:
                r.text = ""
        print(f"  · toc 页{sidx}: 重排为 {len(items)} 条")


def apply_para_replace(prs, spec):
    """按'包含'匹配整个段落并整段改写（用于节标题残留，如 '第3节电路中的电能' → '第4节核能'）。
    只替换含关键词的那个段落，保留该段第一个 run 的格式。
    spec: {"<页号>": [{"contains": "...", "text": "..."}]}
    """
    for sidx, rules in spec.items():
        slide = prs.slides[int(sidx) - 1]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                for rule in rules:
                    if rule["contains"] in full:
                        if not para.runs:
                            para.add_run()
                        para.runs[0].text = rule["text"]
                        for r in para.runs[1:]:
                            r.text = ""
                        break


def apply_box(prs, spec):
    """填充/改写一个文本框（用于『导入新课』右侧内容框填课程引导 等）。
    spec: {"<页号>": {"shape": "形状名(可选，缺省取该页最大/最靠右文本框)",
                      "text": "...", "mode": "replace|append"}}
    """
    for sidx, cfg in spec.items():
        slide = prs.slides[int(sidx) - 1]
        shape = None
        if cfg.get("shape"):
            shape = next((s for s in slide.shapes if s.name == cfg["shape"]), None)
        if shape is None:
            cands = [s for s in slide.shapes
                     if getattr(s, "has_text_frame", False) and s.left is not None]
            if not cands:
                print(f"[warn] box 页{sidx} 无文本框，跳过")
                continue
            shape = max(cands, key=lambda s: (s.width * s.height, s.left))
        if not getattr(shape, "has_text_frame", False):
            print(f"[warn] box 页{sidx} 目标不是文本框，跳过")
            continue
        text = cfg["text"]
        if cfg.get("mode") == "append":
            p = shape.text_frame.add_paragraph()
            if not p.runs:
                p.add_run()
            p.runs[0].text = text
        else:
            replace_text(shape, text)
        print(f"  · box 页{sidx}: 已填充文本框")
def main():
    ap = argparse.ArgumentParser(description="PPTX 排版统一应用")
    ap.add_argument("-i", "--input", required=True, help="输入 .pptx")
    ap.add_argument("-o", "--output", required=True, help="输出 .pptx")
    ap.add_argument("--plan", help="方案 JSON")
    ap.add_argument("--append-teaching", help="教学设计 JSON(单独)或留空用 plan.teaching")
    args = ap.parse_args()

    if not args.input.endswith(".pptx"):
        sys.exit("仅支持 .pptx（.ppt 请先用 Office/LibreOffice 另存为 .pptx）")
    if args.output == args.input:
        sys.exit("输出与输入相同, 会覆盖原件; 请换输出名或确认覆盖意图")

    plan = {}
    if args.plan:
        with open(args.plan, encoding="utf-8") as f:
            plan = json.load(f)

    prs = Presentation(args.input)

    if plan.get("titles"):
        apply_titles(prs, plan["titles"])
    if plan.get("text_replace"):
        apply_text_replace(prs, plan["text_replace"])
    if plan.get("table_replace"):
        apply_table_replace(prs, plan["table_replace"])
    if plan.get("replace_images"):
        apply_replace_images(prs, plan["replace_images"])
    if plan.get("round_rect"):
        apply_round_rect(prs, plan["round_rect"])
    if plan.get("layout_images"):
        apply_layout_images(prs, plan["layout_images"])
    if plan.get("fonts"):
        apply_fonts(prs, plan["fonts"])
    if plan.get("size_map"):
        apply_sizes(prs, plan["size_map"])
    if plan.get("images"):
        apply_images(prs, plan["images"])

    if plan.get("toc"):
        apply_toc(prs, plan["toc"])
    if plan.get("para_replace"):
        apply_para_replace(prs, plan["para_replace"])
    if plan.get("box"):
        apply_box(prs, plan["box"])

    teaching = plan.get("teaching")
    if args.append_teaching:
        with open(args.append_teaching, encoding="utf-8") as f:
            teaching = json.load(f)
    if teaching:
        add_teaching_slide(prs, teaching)

    prs.save(args.output)
    print(f"✓ 已保存: {args.output}")


if __name__ == "__main__":
    main()
